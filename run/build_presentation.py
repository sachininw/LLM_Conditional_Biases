"""
build_presentation.py — Generate experiments section slides as a .pptx file.

Usage:
    python build_presentation.py
Output:
    run/data/experiments_section.pptx
"""

import os
import io
import textwrap
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import seaborn as sns
from scipy import stats

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────────
_RUN_DIR        = os.path.dirname(os.path.abspath(__file__))
_PILOT_CSV      = os.path.join(_RUN_DIR, "data", "conditional_decision_results", "GPT-4o_pilot.csv")
_PLOT_DIR       = os.path.join(_RUN_DIR, "data", "plots")
_EXT_PLOT_DIR   = os.path.join(_RUN_DIR, "data", "plots", "extended")
_MULTI_PLOT_DIR = os.path.join(_RUN_DIR, "data", "plots", "multi_llm")
_OUT_PPTX       = os.path.join(_RUN_DIR, "data", "experiments_section.pptx")

BIAS_ORDER = [
    "Anchoring", "Availability Heuristic", "Bandwagon Effect",
    "Confirmation Bias", "Framing Effect", "In Group Bias",
    "Loss Aversion", "Planning Fallacy", "Status Quo Bias",
]

# ── brand colours ──────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0F, 0x17, 0x2A)   # dark navy background
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT   = RGBColor(0x4A, 0x9E, 0xFF)   # blue accent
C_RED      = RGBColor(0xE5, 0x5A, 0x4E)   # suppress / negative
C_GREEN    = RGBColor(0x4C, 0xC9, 0x8A)   # amplify / positive
C_GREY     = RGBColor(0x8A, 0x93, 0xA8)   # muted
C_YELLOW   = RGBColor(0xF5, 0xC5, 0x42)   # highlight

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, colour=C_BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, colour=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = colour
    return txb


def add_rect(slide, left, top, width, height, fill=C_ACCENT, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_image_from_fig(slide, fig, left, top, width, height):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    slide.shapes.add_picture(buf, left, top, width, height)
    plt.close(fig)


def add_image_file(slide, path, left, top, width, height):
    slide.shapes.add_picture(path, left, top, width, height)


def title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), fill=C_ACCENT)
    add_text(slide, title, Inches(0.35), Inches(0.08), Inches(12.5), Inches(0.65),
             size=28, bold=True, colour=C_BG, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.72), Inches(12.5), Inches(0.35),
                 size=13, bold=False, colour=C_BG, align=PP_ALIGN.LEFT)


def bullet(slide, items, left, top, width, height,
           size=15, colour=C_WHITE, spacing=0.38):
    for i, item in enumerate(items):
        add_text(slide, f"▸  {item}",
                 left, top + Inches(i * spacing), width, Inches(spacing + 0.05),
                 size=size, colour=colour)


def section_tag(slide, label):
    add_rect(slide, Inches(11.8), Inches(0.15), Inches(1.3), Inches(0.4), fill=C_YELLOW)
    add_text(slide, label, Inches(11.82), Inches(0.17), Inches(1.26), Inches(0.38),
             size=10, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)


def _try_img(directory, name):
    path = os.path.join(directory, name)
    return path if os.path.exists(path) else None


def _placeholder_box(slide, left, top, width, height,
                     msg="Run analysis script to generate this plot"):
    add_rect(slide, left, top, width, height, fill=RGBColor(0x1A, 0x23, 0x3D))
    mid_y = top + (height - Inches(0.45)) // 2
    add_rect(slide, left + Inches(0.2), mid_y, width - Inches(0.4), Inches(0.45),
             fill=RGBColor(0x0F, 0x17, 0x2A))
    add_text(slide, msg, left + Inches(0.2), mid_y + Inches(0.04),
             width - Inches(0.4), Inches(0.38),
             size=10, colour=C_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# DATA
# ══════════════════════════════════════════════════════════════════════════════

def load():
    df = pd.read_csv(_PILOT_CSV)
    jp = df[(df["status"] == "OK") & (df["jury_passed"] == True)].copy()
    jp["is_diagonal"] = (jp["bias"] == jp["human_bias"]).astype(int)
    pivot = jp.pivot_table(values="delta_score", index="bias",
                           columns="human_bias", aggfunc="mean")
    pivot = pivot.reindex(index=BIAS_ORDER, columns=BIAS_ORDER)
    return jp, pivot


# ══════════════════════════════════════════════════════════════════════════════
# FIGURE HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def fig_coupling_matrix(jp, pivot):
    fig, ax = plt.subplots(figsize=(9, 6.5))
    fig.patch.set_facecolor("#0F172A")
    ax.set_facecolor("#0F172A")

    mask = pivot.isnull()
    sns.heatmap(
        pivot.round(2), ax=ax,
        cmap=sns.diverging_palette(220, 20, as_cmap=True),
        center=0, annot=True, fmt=".2f", annot_kws={"size": 7.5},
        vmin=-1, vmax=1, mask=mask,
        linewidths=0.4, linecolor="#1E293B",
        cbar_kws={"shrink": 0.6, "label": "Mean Δ-score"},
    )
    # Hatching for missing cells
    for i, rb in enumerate(BIAS_ORDER):
        for j, cb in enumerate(BIAS_ORDER):
            if np.isnan(pivot.loc[rb, cb] if rb in pivot.index and cb in pivot.columns else np.nan):
                ax.add_patch(plt.Rectangle((j, i), 1, 1, fill=True,
                                           color="#1E293B", zorder=2))
    # Bold diagonal
    for k in range(len(BIAS_ORDER)):
        ax.add_patch(plt.Rectangle((k, k), 1, 1, fill=False,
                                   edgecolor="#F5C542", lw=2, zorder=3))

    ax.set_xlabel("Human Bias  →  expressed by interlocutor", color="white", fontsize=9, labelpad=6)
    ax.set_ylabel("Target Bias  →  measured in model", color="white", fontsize=9, labelpad=6)
    ax.tick_params(colors="white", labelsize=7.5)
    ax.set_xticklabels([b.replace(" ", "\n") for b in BIAS_ORDER], rotation=0, ha="center")
    ax.set_yticklabels(BIAS_ORDER, rotation=0)
    ax.xaxis.tick_top(); ax.xaxis.set_label_position("top")
    for spine in ax.spines.values(): spine.set_visible(False)
    cbar = ax.collections[0].colorbar
    cbar.ax.yaxis.label.set_color("white")
    cbar.ax.tick_params(colors="white")

    fig.suptitle("Bias Coupling Matrix  M[β, γ]  (jury-confirmed, N=273)",
                 color="white", fontsize=11, y=0.02)
    fig.tight_layout()
    return fig


def fig_sycophancy_test(pivot):
    cells = pivot.values.flatten()
    cells = cells[~np.isnan(cells)]
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    fig.patch.set_facecolor("#0F172A"); ax.set_facecolor("#151F35")

    pos = cells[cells > 0]; neg = cells[cells < 0]; zer = cells[cells == 0]
    ax.scatter(pos, np.random.uniform(0.3, 0.7, len(pos)), color="#4CC98A", s=50, alpha=0.85, label=f"Positive ({len(pos)})", zorder=3)
    ax.scatter(neg, np.random.uniform(0.3, 0.7, len(neg)), color="#E55A4E", s=50, alpha=0.85, label=f"Negative ({len(neg)})", zorder=3)
    ax.scatter(zer, np.random.uniform(0.3, 0.7, len(zer)), color="#8A93A8", s=50, alpha=0.6, label=f"Zero ({len(zer)})", zorder=3)
    ax.axvline(0, color="white", lw=1.2, linestyle="--")
    ax.axvspan(0, 1.05, alpha=0.08, color="#4CC98A", label="Sycophancy predicts: all here")
    ax.set_xlim(-1.1, 1.1); ax.set_ylim(0, 1)
    ax.set_xlabel("Cell mean Δ-score", color="white", fontsize=9)
    ax.set_yticks([]); ax.tick_params(colors="white")
    for sp in ax.spines.values(): sp.set_color("#4A9EFF")
    ax.legend(fontsize=7.5, facecolor="#0F172A", edgecolor="#4A9EFF",
              labelcolor="white", loc="upper left")
    ax.set_title("Sycophancy test: cell sign distribution", color="white", fontsize=9, pad=6)
    fig.tight_layout()
    return fig


def fig_style_mirroring_test(pivot):
    biases_with_diag = [b for b in BIAS_ORDER
                        if b in pivot.index and b in pivot.columns
                        and not np.isnan(pivot.loc[b, b])]
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    fig.patch.set_facecolor("#0F172A"); ax.set_facecolor("#151F35")

    y_pos = list(range(len(biases_with_diag)))
    for yi, bias in enumerate(biases_with_diag):
        row = pivot.loc[bias].dropna()
        off = row.drop(bias, errors="ignore").values
        diag_v = row[bias]
        rank = int((row > diag_v).sum()) + 1
        total = len(row)
        # off-diagonal as grey dots
        ax.scatter(off, [yi] * len(off), color="#8A93A8", s=30, alpha=0.7, zorder=2)
        # diagonal as red/green dot
        color = "#4CC98A" if diag_v > 0 else "#E55A4E"
        ax.scatter([diag_v], [yi], color=color, s=90, zorder=4,
                   edgecolors="#F5C542", linewidths=1.5)
        ax.text(1.05, yi, f"rank {rank}/{total}", va="center",
                fontsize=7.5, color="#F5C542")

    ax.axvline(0, color="white", lw=1, linestyle="--")
    ax.set_yticks(y_pos)
    ax.set_yticklabels([b.replace(" ", "\n") for b in biases_with_diag], fontsize=7)
    ax.tick_params(colors="white")
    ax.set_xlabel("Δ-score", color="white", fontsize=9)
    for sp in ax.spines.values(): sp.set_color("#4A9EFF")
    ax.set_title("Style mirroring test: diagonal rank in its row\n(yellow dot = diagonal, grey = off-diagonal)",
                 color="white", fontsize=8.5, pad=5)
    ax.set_xlim(-1.05, 1.35)
    fig.tight_layout()
    return fig


def fig_row_col_structure(jp):
    row_stats = jp.groupby("bias")["delta_score"].agg(["mean", "sem"]).reindex(BIAS_ORDER)
    col_stats = jp.groupby("human_bias")["delta_score"].agg(["mean", "sem"]).reindex(BIAS_ORDER)

    fig, axes = plt.subplots(1, 2, figsize=(10, 4.2))
    fig.patch.set_facecolor("#0F172A")

    for ax, stats_df, label in [
        (axes[0], row_stats, "Target susceptibility\n(row mean Δ)"),
        (axes[1], col_stats,  "Primer strength\n(column mean Δ)"),
    ]:
        ax.set_facecolor("#151F35")
        ordered = stats_df.sort_values("mean")
        colors = []
        for b, row in ordered.iterrows():
            m = row["mean"]
            if b == "In Group Bias" and label.startswith("Target"):
                colors.append("#E55A4E")
            elif m > 0.05:
                colors.append("#4CC98A")
            elif m < -0.05:
                colors.append("#E55A4E")
            else:
                colors.append("#8A93A8")
        y = np.arange(len(ordered))
        ax.barh(y, ordered["mean"], xerr=ordered["sem"],
                color=colors, edgecolor="none", height=0.6,
                error_kw={"ecolor": "white", "capsize": 3, "lw": 1})
        ax.axvline(0, color="white", lw=1, linestyle="--")
        ax.set_yticks(y)
        ax.set_yticklabels(ordered.index, fontsize=8, color="white")
        ax.tick_params(colors="white", labelsize=8)
        ax.set_xlabel("Mean Δ-score  ±SE", color="white", fontsize=8)
        ax.set_title(label, color="white", fontsize=9, pad=6)
        for sp in ax.spines.values(): sp.set_color("#4A9EFF")

    fig.suptitle("Which biases are susceptible targets? Which are effective primers?",
                 color="white", fontsize=10, y=1.01)
    fig.tight_layout()
    return fig


def fig_significance_summary(jp):
    from statsmodels.stats.multitest import multipletests
    rows = []
    for bias in BIAS_ORDER:
        vals = jp[jp["bias"] == bias]["delta_score"].dropna()
        if len(vals) < 2:
            rows.append({"bias": bias, "mean": np.nan, "d": np.nan, "p": np.nan})
            continue
        t_, p_ = stats.ttest_1samp(vals, 0)
        d_ = vals.mean() / (vals.std() + 1e-9)
        rows.append({"bias": bias, "n": len(vals), "mean": vals.mean(), "d": d_, "p": p_})
    df2 = pd.DataFrame(rows).set_index("bias")

    fig, axes = plt.subplots(1, 2, figsize=(10, 4))
    fig.patch.set_facecolor("#0F172A")

    # Left: Cohen's d lollipop
    ax = axes[0]; ax.set_facecolor("#151F35")
    ordered = df2["d"].reindex(BIAS_ORDER).fillna(0)
    colors = ["#E55A4E" if v < -0.5 else "#4CC98A" if v > 0.5 else "#8A93A8" for v in ordered]
    y = np.arange(len(BIAS_ORDER))
    ax.barh(y, ordered, color=colors, height=0.5, edgecolor="none")
    ax.scatter(ordered, y, color=colors, s=60, zorder=4)
    ax.axvline(0, color="white", lw=1, ls="--")
    ax.axvline( 0.5, color="#4CC98A", lw=0.8, ls=":", alpha=0.6)
    ax.axvline(-0.5, color="#E55A4E", lw=0.8, ls=":", alpha=0.6)
    ax.set_yticks(y); ax.set_yticklabels(BIAS_ORDER, fontsize=8, color="white")
    ax.tick_params(colors="white"); ax.set_xlabel("Cohen's d", color="white", fontsize=9)
    ax.set_title("Effect size per target bias", color="white", fontsize=9, pad=5)
    for sp in ax.spines.values(): sp.set_color("#4A9EFF")

    # Right: -log10(p) bar
    ax2 = axes[1]; ax2.set_facecolor("#151F35")
    pvals = df2["p"].reindex(BIAS_ORDER)
    logp  = -np.log10(pvals.clip(lower=1e-6))
    colors2 = ["#E55A4E" if (p is not None and not np.isnan(p) and p < 0.05) else "#8A93A8"
               for p in pvals]
    ax2.barh(y, logp, color=colors2, height=0.5, edgecolor="none")
    ax2.axvline(-np.log10(0.05), color="#F5C542", lw=1.2, ls="--", label="p=0.05")
    ax2.set_yticks(y); ax2.set_yticklabels(BIAS_ORDER, fontsize=8, color="white")
    ax2.tick_params(colors="white")
    ax2.set_xlabel("−log₁₀(p)", color="white", fontsize=9)
    ax2.set_title("Significance per target bias", color="white", fontsize=9, pad=5)
    ax2.legend(fontsize=8, facecolor="#0F172A", edgecolor="#4A9EFF", labelcolor="white")
    for sp in ax2.spines.values(): sp.set_color("#4A9EFF")

    fig.tight_layout()
    return fig


def fig_causal(jp):
    from scipy.optimize import minimize
    fig, axes = plt.subplots(1, 2, figsize=(10, 4))
    fig.patch.set_facecolor("#0F172A")

    # Left: DiD parallel trends for In Group Bias
    ax = axes[0]; ax.set_facecolor("#151F35")
    igb = jp[jp["bias"] == "In Group Bias"]
    ind_mean = igb["individual_score"].mean()
    neu_mean = igb["neutral_score"].mean()
    ax.plot([0, 1], [0, ind_mean], "o-", color="#E55A4E", lw=2, ms=8, label=f"Biased human (post={ind_mean:.3f})")
    ax.plot([0, 1], [0, neu_mean], "s--", color="#4A9EFF", lw=2, ms=8, label=f"Neutral human (post={neu_mean:.3f})")
    ax.axhline(0, color="white", lw=0.8, ls=":")
    att = ind_mean - neu_mean
    ax.annotate(f"ATT = {att:.3f}", xy=(1, (ind_mean + neu_mean) / 2),
                xytext=(0.6, (ind_mean + neu_mean) / 2 + 0.05),
                fontsize=9, color="#F5C542",
                arrowprops=dict(arrowstyle="->", color="#F5C542"))
    ax.set_xticks([0, 1]); ax.set_xticklabels(["Pre\n(no message)", "Post\n(w/ message)"], color="white")
    ax.tick_params(colors="white"); ax.set_ylabel("Score", color="white", fontsize=9)
    ax.set_title("DiD: In-Group Bias  (jury-confirmed)", color="white", fontsize=9, pad=5)
    ax.legend(fontsize=8, facecolor="#0F172A", edgecolor="#4A9EFF", labelcolor="white")
    for sp in ax.spines.values(): sp.set_color("#4A9EFF")

    # Right: SCM gap bars
    ax2 = axes[1]; ax2.set_facecolor("#151F35")
    scm_data = {
        "Anchoring": (-0.250, 0.000),
        "Availability Heuristic": (0.037, 0.014),
        "Bandwagon Effect": (0.000, 0.000),
        "Confirmation Bias": (0.000, 0.000),
        "In Group Bias": (-0.400, -0.257),
        "Loss Aversion": (-0.011, 0.146),
        "Status Quo Bias": (0.150, -0.120),
    }
    p_vals = {"Anchoring": 0.000, "Availability Heuristic": 1.000,
              "Bandwagon Effect": 1.000, "Confirmation Bias": 1.000,
              "In Group Bias": 0.500, "Loss Aversion": 0.333, "Status Quo Bias": 0.000}
    gaps = {b: a - s for b, (a, s) in scm_data.items()}
    biases_scm = list(gaps.keys())
    gap_vals = [gaps[b] for b in biases_scm]
    colors_scm = ["#E55A4E" if p_vals[b] < 0.05 else "#F5C542" if p_vals[b] < 0.10 else "#8A93A8"
                  for b in biases_scm]
    y2 = np.arange(len(biases_scm))
    ax2.barh(y2, gap_vals, color=colors_scm, height=0.55, edgecolor="none")
    ax2.axvline(0, color="white", lw=1, ls="--")
    ax2.set_yticks(y2); ax2.set_yticklabels(biases_scm, fontsize=8, color="white")
    ax2.tick_params(colors="white")
    ax2.set_xlabel("SCM Gap  (actual − synthetic Δ)", color="white", fontsize=9)
    ax2.set_title("Synthetic Control: diagonal gaps", color="white", fontsize=9, pad=5)
    for sp in ax2.spines.values(): sp.set_color("#4A9EFF")
    legend_els = [mpatches.Patch(color="#E55A4E", label="p<0.05"),
                  mpatches.Patch(color="#F5C542", label="p<0.10"),
                  mpatches.Patch(color="#8A93A8", label="n.s.")]
    ax2.legend(handles=legend_els, fontsize=8, facecolor="#0F172A",
               edgecolor="#4A9EFF", labelcolor="white")
    fig.tight_layout()
    return fig


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = blank_slide(prs); fill_bg(sl)
    # Large centre title
    add_rect(sl, Inches(0), Inches(2.8), SLIDE_W, Inches(1.6), fill=C_ACCENT)
    add_text(sl,
             "Experiments: Conditional Cognitive Bias in LLMs",
             Inches(0.5), Inches(2.9), Inches(12.3), Inches(1.2),
             size=34, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
    add_text(sl,
             "Is cognitive bias an intrinsic property of the model,\n"
             "or a relational property of the human–LLM dyad?",
             Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.0),
             size=18, colour=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "GPT-4o  ·  10 Cognitive Biases  ·  N = 273 jury-confirmed interactions",
             Inches(2), Inches(5.8), Inches(9.3), Inches(0.5),
             size=13, colour=C_GREY, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "Section Overview")
    items = [
        ("§4.1",  "Hypothesis & the coupling matrix"),
        ("§4.2",  "Experimental conditions  (Control · Treatment · Neutral)"),
        ("§4.3",  "Dataset & scenario structure"),
        ("§4.4",  "BiasJury — multi-provider validation framework"),
        ("§4.5",  "Scoring metric & Δ computation"),
        ("§4.6",  "Three competing mechanisms + discriminating tests"),
        ("§4.7",  "Statistical framework  (t-tests, FDR, LMM)"),
        ("§4.8",  "Causal validation  (DiD · PSM · SCM)"),
        ("§4.9",  "Extended analyses  (semantic · bootstrap · SVD · taxonomy · regression · network)"),
        ("§4.10", "Multi-LLM replication study  (5 models, cross-model coupling comparison)"),
    ]
    spacing = 0.60
    for i, (tag, desc) in enumerate(items):
        top = Inches(1.25 + i * spacing)
        add_rect(sl, Inches(0.4), top, Inches(0.9), Inches(0.48), fill=C_ACCENT)
        add_text(sl, tag, Inches(0.4), top + Inches(0.03), Inches(0.9), Inches(0.42),
                 size=11, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
        add_text(sl, desc, Inches(1.45), top + Inches(0.03), Inches(10.5), Inches(0.42),
                 size=13, colour=C_WHITE)


def slide_hypothesis(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.1  Core Hypothesis", "What are we actually testing?")
    section_tag(sl, "§4.1")

    # Main hypothesis box
    add_rect(sl, Inches(0.4), Inches(1.25), Inches(12.5), Inches(1.3), fill=RGBColor(0x1E,0x29,0x4A))
    add_text(sl,
             "Cognitive bias in deployed LLMs is not an intrinsic property of the model\n"
             "but a relational property of the human–LLM dyad.",
             Inches(0.55), Inches(1.32), Inches(12.2), Inches(1.1),
             size=17, bold=True, colour=C_YELLOW, align=PP_ALIGN.CENTER)

    add_text(sl, "Operationalised as the Bias Coupling Matrix  M[β, γ]",
             Inches(0.4), Inches(2.72), Inches(12.5), Inches(0.4),
             size=14, bold=True, colour=C_ACCENT)

    add_text(sl,
             "M[β, γ]  =  E[ Δ  |  target = β,  human expressed = γ ]",
             Inches(1.0), Inches(3.15), Inches(11.0), Inches(0.45),
             size=15, bold=True, colour=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "where  Δ = individual_score − neutral_score  isolates bias-specific induction from any human-turn baseline effect",
             Inches(1.0), Inches(3.65), Inches(11.0), Inches(0.4),
             size=12, colour=C_GREY, align=PP_ALIGN.CENTER)

    # Three outcomes
    for x, col, title, body in [
        (Inches(0.35), C_RED,
         "If M is diagonal & degenerate",
         "Collapses to amplification / sycophancy.\nNothing new established."),
        (Inches(4.6), C_GREY,
         "If M is uniformly positive",
         "Pure sycophancy: model agrees with\nwhatever human says. Not novel."),
        (Inches(8.85), C_GREEN,
         "If M is structured & asymmetric",
         "Mechanism-revealing. Evidence for\nthe dyadic hypothesis. Genuinely new."),
    ]:
        add_rect(sl, x, Inches(4.25), Inches(3.9), Inches(2.85),
                 fill=RGBColor(0x1E,0x29,0x4A))
        add_rect(sl, x, Inches(4.25), Inches(3.9), Inches(0.35), fill=col)
        add_text(sl, title, x + Inches(0.1), Inches(4.28),
                 Inches(3.7), Inches(0.3), size=11, bold=True, colour=C_BG)
        add_text(sl, body, x + Inches(0.1), Inches(4.65),
                 Inches(3.7), Inches(2.3), size=12, colour=C_WHITE)


def slide_three_mechanisms(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.6  Three Competing Mechanisms", "The coupling matrix discriminates between them")
    section_tag(sl, "§4.6")

    mechs = [
        (C_RED, "1  Pure Sycophancy",
         "Model agrees with whatever the human says, regardless of bias type.",
         "Predicts: M uniformly positive, diagonal ≈ off-diagonal, CV ≈ 1"),
        (C_ACCENT, "2  Style Mirroring",
         "Model mirrors the human's cognitive frame within the same bias family only.",
         "Predicts: diagonal dominance — M[β,β] >> M[β,γ] for β≠γ"),
        (C_GREEN, "3  Latent Priming",
         "Human message activates representational neighbourhood; cross-induction follows semantic topology.",
         "Predicts: structured off-diagonal, asymmetric M, correlated with semantic similarity s(β,γ)"),
    ]
    for i, (col, name, desc, pred) in enumerate(mechs):
        top = Inches(1.3 + i * 1.85)
        add_rect(sl, Inches(0.35), top, Inches(0.18), Inches(1.6), fill=col)
        add_rect(sl, Inches(0.6), top, Inches(12.3), Inches(1.6), fill=RGBColor(0x1A,0x23,0x3D))
        add_text(sl, name, Inches(0.75), top + Inches(0.1), Inches(11.8), Inches(0.4),
                 size=15, bold=True, colour=col)
        add_text(sl, desc, Inches(0.75), top + Inches(0.48), Inches(11.8), Inches(0.45),
                 size=13, colour=C_WHITE)
        add_text(sl, pred, Inches(0.75), top + Inches(0.95), Inches(11.8), Inches(0.5),
                 size=12, colour=C_GREY)

    add_text(sl, "Our pilot data: CV = 10.6  ·  20 positive / 23 negative / 26 zero cells  ·  diagonal ranks last in 3/7 rows",
             Inches(0.35), Inches(6.85), Inches(12.6), Inches(0.4),
             size=11, colour=C_YELLOW, align=PP_ALIGN.CENTER)


def slide_protocol(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.2  Experimental Protocol", "Three conditions — one experimental object")
    section_tag(sl, "§4.2")

    cond_data = [
        (Inches(0.35),  C_GREY,  "CONTROL",
         ["Single turn only", "No human message", "Establishes baseline decision"],
         ["Diverging situation + options  →  LLM answer"]),
        (Inches(4.6),   C_ACCENT, "NEUTRAL",
         ["Four-turn conversation", "Jury-validated neutral human", "Controls for human-turn effect"],
         ["Common situation  →  Neutral message (jury ✓)  →  LLM response  →  Decision"]),
        (Inches(8.85),  C_GREEN, "TREATMENT",
         ["Four-turn conversation", "Jury-confirmed biased human", "Measures conditional bias"],
         ["Common situation  →  Biased message (jury ✓)  →  LLM response  →  Decision"]),
    ]
    for x, col, label, bullets, flow in cond_data:
        add_rect(sl, x, Inches(1.25), Inches(4.0), Inches(5.9), fill=RGBColor(0x1A,0x23,0x3D))
        add_rect(sl, x, Inches(1.25), Inches(4.0), Inches(0.42), fill=col)
        add_text(sl, label, x + Inches(0.1), Inches(1.28), Inches(3.8), Inches(0.35),
                 size=13, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
        for bi, b in enumerate(bullets):
            add_text(sl, f"• {b}", x + Inches(0.15), Inches(1.75 + bi * 0.42),
                     Inches(3.7), Inches(0.38), size=12, colour=C_WHITE)
        add_rect(sl, x + Inches(0.1), Inches(3.65), Inches(3.8), Inches(0.03), fill=C_GREY)
        add_text(sl, flow[0], x + Inches(0.15), Inches(3.75), Inches(3.7), Inches(2.9),
                 size=10, colour=C_GREY)

    # Delta formula at bottom
    add_rect(sl, Inches(0.35), Inches(6.78), Inches(12.6), Inches(0.45), fill=RGBColor(0x1E,0x29,0x4A))
    add_text(sl,
             "Δ  =  individual_score  −  neutral_score     →     M[β,γ]  =  E[Δ]",
             Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.4),
             size=14, bold=True, colour=C_YELLOW, align=PP_ALIGN.CENTER)


def slide_jury(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.4  BiasJury — Multi-Provider Validation", "Eliminates evaluator circularity")
    section_tag(sl, "§4.4")

    # Three juror cards
    jurors = [
        ("J1", "GPT-4-Turbo", "OpenAI", C_ACCENT),
        ("J2", "Claude Haiku 4.5", "Anthropic", C_GREEN),
        ("J3", "GPT-4o-Mini", "OpenAI", C_YELLOW),
    ]
    for i, (jid, model, provider, col) in enumerate(jurors):
        x = Inches(0.35 + i * 3.0)
        add_rect(sl, x, Inches(1.25), Inches(2.7), Inches(1.7), fill=RGBColor(0x1A,0x23,0x3D))
        add_rect(sl, x, Inches(1.25), Inches(2.7), Inches(0.38), fill=col)
        add_text(sl, jid, x + Inches(0.1), Inches(1.27), Inches(2.5), Inches(0.32),
                 size=14, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
        add_text(sl, model, x + Inches(0.1), Inches(1.7), Inches(2.5), Inches(0.4),
                 size=12, bold=True, colour=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, provider, x + Inches(0.1), Inches(2.1), Inches(2.5), Inches(0.35),
                 size=11, colour=C_GREY, align=PP_ALIGN.CENTER)

    # Tasks
    add_text(sl, "Each juror performs independently:", Inches(9.7), Inches(1.25), Inches(3.5), Inches(0.35),
             size=12, bold=True, colour=C_ACCENT)
    for bi, b in enumerate(["Rating pass: score all 30 biases 0–5 + naturalness",
                             "Blind ID pass: name top-3 biases without knowing target"]):
        add_text(sl, f"• {b}", Inches(9.7), Inches(1.65 + bi * 0.45), Inches(3.45), Inches(0.4),
                 size=11, colour=C_WHITE)

    # 4 criteria
    add_text(sl, "A message PASSES all four criteria jointly:",
             Inches(0.35), Inches(3.1), Inches(8.5), Inches(0.38),
             size=13, bold=True, colour=C_WHITE)

    criteria = [
        ("Mean target-bias rating", "≥ 3.0", "pilot mean: 3.99"),
        ("Purity  (target − max other)", "≥ 1.5", "pilot mean: 1.82"),
        ("Mean naturalness", "≥ 3.0", "pilot mean: 3.97"),
        ("Blind ID accuracy  (top-1 correct)", "≥ 2/3", "Krippendorff α = 0.42"),
    ]
    for i, (crit, thresh, stat) in enumerate(criteria):
        row_top = Inches(3.55 + i * 0.72)
        add_rect(sl, Inches(0.35), row_top, Inches(5.5), Inches(0.6), fill=RGBColor(0x1A,0x23,0x3D))
        add_text(sl, crit, Inches(0.5), row_top + Inches(0.08), Inches(4.0), Inches(0.45), size=12, colour=C_WHITE)
        add_rect(sl, Inches(5.95), row_top, Inches(1.1), Inches(0.6), fill=C_ACCENT)
        add_text(sl, thresh, Inches(5.97), row_top + Inches(0.08), Inches(1.06), Inches(0.45),
                 size=12, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
        add_text(sl, stat, Inches(7.2), row_top + Inches(0.1), Inches(3.0), Inches(0.42),
                 size=11, colour=C_GREY)

    # Pass rates
    add_rect(sl, Inches(0.35), Inches(6.52), Inches(9.0), Inches(0.65), fill=RGBColor(0x1E,0x29,0x4A))
    add_text(sl,
             "Pilot pass rate:  Biased messages 55.3%  (273/494)     ·     Neutral messages 100%  (494/494)",
             Inches(0.5), Inches(6.6), Inches(8.7), Inches(0.4),
             size=13, colour=C_YELLOW)


def slide_metric(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.5  Scoring Metric & Δ Computation")
    section_tag(sl, "§4.5")

    add_text(sl, "RatioScaleMetric  (from the original benchmark):",
             Inches(0.4), Inches(1.3), Inches(12.0), Inches(0.38),
             size=14, bold=True, colour=C_ACCENT)

    add_rect(sl, Inches(1.5), Inches(1.75), Inches(10.0), Inches(0.85), fill=RGBColor(0x1A,0x23,0x3D))
    add_text(sl,
             "𝔅  =  k · ( |â_ctrl − x₁|  −  |â_treat − x₂| )  /  max( |â_ctrl − x₁|, |â_treat − x₂| )",
             Inches(1.6), Inches(1.8), Inches(9.8), Inches(0.72),
             size=16, bold=True, colour=C_WHITE, align=PP_ALIGN.CENTER)

    params = [
        ("k ∈ {−1, +1}", "bias direction, from dataset metric_params"),
        ("x₁, x₂", "bias-neutral anchor values"),
        ("â_ctrl, â_treat", "ordinal decision indices under control / treatment"),
        ("Range: [−1, 1]", "+1 = maximal bias,  −1 = maximal counter-bias,  0 = neutral"),
    ]
    for i, (sym, desc) in enumerate(params):
        add_rect(sl, Inches(0.4), Inches(2.72 + i * 0.52), Inches(2.1), Inches(0.44), fill=C_ACCENT)
        add_text(sl, sym, Inches(0.45), Inches(2.75 + i * 0.52), Inches(2.0), Inches(0.38),
                 size=12, bold=True, colour=C_BG)
        add_text(sl, desc, Inches(2.65), Inches(2.75 + i * 0.52), Inches(9.8), Inches(0.38),
                 size=12, colour=C_WHITE)

    add_text(sl, "Three scores per observation:", Inches(0.4), Inches(4.9), Inches(12.0), Inches(0.38),
             size=14, bold=True, colour=C_ACCENT)

    scores = [
        (C_RED,    "individual_score", "𝔅( â_ctrl,  â_biased-treat )", "biased human condition"),
        (C_ACCENT, "neutral_score",    "𝔅( â_ctrl,  â_neutral-treat )", "neutral human baseline"),
        (C_YELLOW, "Δ  (delta)",       "individual_score  −  neutral_score", "isolates bias-specific induction"),
    ]
    for i, (col, name, formula, meaning) in enumerate(scores):
        x = Inches(0.35 + i * 4.32)
        add_rect(sl, x, Inches(5.35), Inches(4.0), Inches(1.8), fill=RGBColor(0x1A,0x23,0x3D))
        add_rect(sl, x, Inches(5.35), Inches(4.0), Inches(0.38), fill=col)
        add_text(sl, name, x + Inches(0.1), Inches(5.37), Inches(3.8), Inches(0.32),
                 size=12, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
        add_text(sl, formula, x + Inches(0.1), Inches(5.8), Inches(3.8), Inches(0.45),
                 size=11, colour=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, meaning, x + Inches(0.1), Inches(6.3), Inches(3.8), Inches(0.65),
                 size=11, colour=C_GREY, align=PP_ALIGN.CENTER)


def slide_matrix(prs, jp, pivot):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "The Bias Coupling Matrix  M[β, γ]",
              "jury_passed=True rows only  ·  N=273  ·  Δ = individual_score − neutral_score")

    fig = fig_coupling_matrix(jp, pivot)
    add_image_from_fig(sl, fig, Inches(0.3), Inches(1.15), Inches(12.7), Inches(6.1))


def slide_mechanism_tests(prs, jp, pivot):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.6  Mechanism Discrimination Tests",
              "One panel per mechanism — each shows what the null predicts vs. what we observe")
    section_tag(sl, "§4.6")

    fig_syco  = fig_sycophancy_test(pivot)
    fig_mirr  = fig_style_mirroring_test(pivot)

    add_image_from_fig(sl, fig_syco, Inches(0.3),  Inches(1.2), Inches(5.9), Inches(3.5))
    add_image_from_fig(sl, fig_mirr, Inches(6.4),  Inches(1.2), Inches(6.7), Inches(3.5))

    # Verdict boxes
    verdicts = [
        (Inches(0.3),  C_RED,   "Sycophancy",
         "REJECTED — 23 negative cells, CV=10.6×, global mean p=0.67"),
        (Inches(6.4), C_RED,   "Style Mirroring",
         "REJECTED — diagonal ranks last or 2nd-last in its own row for 3/7 biases"),
        (Inches(4.55), C_GREEN, "Latent Priming",
         "SUPPORTED — structured, asymmetric, bias-specific suppressor (In-Group Bias, p=0.0002)"),
    ]
    for x, col, label, verdict in verdicts:
        w = Inches(3.85) if x != Inches(4.55) else Inches(4.5)
        top = Inches(4.85)
        add_rect(sl, x, top, w, Inches(0.32), fill=col)
        add_text(sl, label, x + Inches(0.08), top + Inches(0.02), w - Inches(0.15), Inches(0.28),
                 size=11, bold=True, colour=C_BG)
        add_rect(sl, x, top + Inches(0.32), w, Inches(0.8), fill=RGBColor(0x1A,0x23,0x3D))
        add_text(sl, verdict, x + Inches(0.08), top + Inches(0.36), w - Inches(0.15), Inches(0.75),
                 size=10, colour=C_WHITE)

    # Latent priming note
    add_rect(sl, Inches(0.3), Inches(5.82), Inches(12.7), Inches(1.45), fill=RGBColor(0x1E,0x29,0x4A))
    add_text(sl, "Latent Priming — additional evidence:",
             Inches(0.5), Inches(5.88), Inches(12.3), Inches(0.32),
             size=12, bold=True, colour=C_GREEN)
    add_text(sl,
             "• Asymmetric pairs: Anchoring→Avail.Heuristic Δ=+0.40, reverse Δ=−0.09  |  "
             "Conf.Bias→In-Group Δ=0.00, reverse Δ=−0.75\n"
             "• Next step: Spearman ρ between semantic similarity matrix S[β,γ] and M[β,γ] "
             "(off-diagonal) — requires one embeddings API call",
             Inches(0.5), Inches(6.22), Inches(12.3), Inches(0.95),
             size=11, colour=C_WHITE)


def slide_row_col(prs, jp):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "Row & Column Structure of M",
              "Which biases are susceptible targets? Which are effective primers?")
    fig = fig_row_col_structure(jp)
    add_image_from_fig(sl, fig, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))
    add_rect(sl, Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.45), fill=RGBColor(0x1E,0x29,0x4A))
    add_text(sl,
             "In-Group Bias: only significant target (p=0.0002, d=−0.885) — "
             "biased interlocutors trigger corrective response, not amplification. "
             "Bandwagon & Confirmation Bias: zero variance — GPT-4o does not exhibit these biases in these scenarios.",
             Inches(0.5), Inches(6.83), Inches(12.3), Inches(0.4),
             size=10.5, colour=C_YELLOW)


def slide_significance(prs, jp):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.7  Statistical Significance", "Cell-level tests · LMM · FDR correction")
    section_tag(sl, "§4.7")

    fig = fig_significance_summary(jp)
    add_image_from_fig(sl, fig, Inches(0.3), Inches(1.2), Inches(9.5), Inches(5.0))

    # Table summary on right
    table_data = [
        ("In-Group Bias",    "24", "−0.319", "−0.885", "0.0002", "sig."),
        ("Anchoring",        "28", "+0.143", "+0.242", "0.212",  "n.s."),
        ("Status Quo Bias",  "21", "+0.067", "+0.311", "0.170",  "n.s."),
        ("Bandwagon Effect", "29", "0.000",  "0.000",  "—",      "absent"),
        ("Conf. Bias",       "32", "0.000",  "0.000",  "—",      "absent"),
    ]
    headers = ["Bias", "N", "Mean Δ", "d", "p", ""]
    col_widths = [Inches(1.7), Inches(0.4), Inches(0.7), Inches(0.55), Inches(0.65), Inches(0.65)]
    xs = [Inches(9.9), Inches(11.6), Inches(12.0), Inches(12.7), Inches(13.25), Inches(13.9)]
    # re-lay within slide width
    col_widths_in = [1.55, 0.38, 0.65, 0.52, 0.62, 0.6]
    col_widths = [Inches(w) for w in col_widths_in]
    xs = [Inches(9.95 + sum(col_widths_in[:i])) for i in range(len(col_widths_in))]
    # header row
    add_rect(sl, Inches(9.9), Inches(1.22), Inches(3.1), Inches(0.38), fill=C_ACCENT)
    for j, h in enumerate(headers):
        add_text(sl, h, xs[j], Inches(1.25), col_widths[j], Inches(0.33),
                 size=10, bold=True, colour=C_BG)
    for ri, row in enumerate(table_data):
        bg = RGBColor(0x1A,0x23,0x3D) if ri % 2 == 0 else RGBColor(0x15,0x1E,0x35)
        add_rect(sl, Inches(9.9), Inches(1.65 + ri * 0.48), Inches(3.1), Inches(0.46), fill=bg)
        for j, val in enumerate(row):
            col = C_RED if val == "sig." else C_GREY if val in ("absent","n.s.") else C_WHITE
            add_text(sl, val, xs[j], Inches(1.67 + ri * 0.48), col_widths[j], Inches(0.42),
                     size=9.5, colour=col)

    add_text(sl, "Full 10×10 cell-level tests: 4 raw p<0.05, 1 FDR-significant\nLMM: Status Quo Bias p=0.015",
             Inches(9.9), Inches(4.2), Inches(3.1), Inches(0.85),
             size=10, colour=C_GREY)
    add_text(sl, "Power note: mean 3.96 obs/cell after jury filter.\n10+ scenarios/bias needed for cell-level power.",
             Inches(9.9), Inches(5.1), Inches(3.1), Inches(0.8),
             size=10, colour=C_YELLOW)


def slide_causal(prs, jp):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.8  Causal Validation", "DiD  ·  Propensity Score Matching  ·  Synthetic Control")
    section_tag(sl, "§4.8")

    fig = fig_causal(jp)
    add_image_from_fig(sl, fig, Inches(0.3), Inches(1.2), Inches(12.7), Inches(4.5))

    # Method summary strip
    methods = [
        ("DiD", "Biased vs neutral groups × pre/post message period.\nCluster-robust SE at scenario level.\nATT = −0.007, p=0.762 (global)", C_ACCENT),
        ("PSM", "Diagonal vs off-diagonal matched on sim. features.\n1:1 nearest-neighbour, sim features only.\nMatched ATT = −0.061, p=0.41", C_YELLOW),
        ("SCM", "Synthetic counterfactual from donor off-diagonal cells.\nPermutation placebo p-values.\nAnchoring & Status Quo Bias: p<0.05", C_GREEN),
    ]
    for i, (name, desc, col) in enumerate(methods):
        x = Inches(0.35 + i * 4.32)
        add_rect(sl, x, Inches(5.85), Inches(4.0), Inches(1.4), fill=RGBColor(0x1A,0x23,0x3D))
        add_rect(sl, x, Inches(5.85), Inches(4.0), Inches(0.32), fill=col)
        add_text(sl, name, x + Inches(0.1), Inches(5.87), Inches(3.8), Inches(0.28),
                 size=12, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
        add_text(sl, desc, x + Inches(0.1), Inches(6.22), Inches(3.8), Inches(0.95),
                 size=10, colour=C_WHITE)


def slide_summary(prs):
    sl = blank_slide(prs); fill_bg(sl)
    add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(1.4), fill=C_ACCENT)
    add_text(sl, "Summary: What the Experiments Establish",
             Inches(0.4), Inches(0.15), Inches(12.5), Inches(1.1),
             size=26, bold=True, colour=C_BG)

    findings = [
        (C_RED,    "Sycophancy rejected",
         "Cell mean CV=10.6×, 23/69 cells negative, 26/69 zero. The coupling matrix is not uniformly positive."),
        (C_RED,    "Style mirroring rejected",
         "Diagonal ranks last or 2nd-last in its own row for 3/7 biases. Same-bias human does NOT produce the highest delta."),
        (C_GREEN,  "In-Group Bias: robust suppressor",
         "Row mean Δ=−0.319, p=0.0002, d=−0.885. Biased interlocutors trigger corrective response across all human bias types."),
        (C_ACCENT, "Matrix is asymmetric & structured",
         "Anchoring←Avail.Heuristic Δ=+0.40 vs reverse Δ=−0.09. Structured off-diagonal — consistent with latent priming."),
        (C_YELLOW, "Power limitation",
         "Mean 3.96 obs/cell after jury filter (55.3% pass rate). 10+ scenarios/bias needed for cell-level inference."),
    ]
    for i, (col, title, body) in enumerate(findings):
        top = Inches(1.5 + i * 1.12)
        add_rect(sl, Inches(0.35), top, Inches(0.15), Inches(0.95), fill=col)
        add_rect(sl, Inches(0.6), top, Inches(12.3), Inches(0.95), fill=RGBColor(0x1A,0x23,0x3D))
        add_text(sl, title, Inches(0.75), top + Inches(0.06), Inches(12.0), Inches(0.35),
                 size=13, bold=True, colour=col)
        add_text(sl, body,  Inches(0.75), top + Inches(0.44), Inches(12.0), Inches(0.45),
                 size=11.5, colour=C_WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# EXTENDED ANALYSIS SLIDES  (§4.9 — §4.10)
# ══════════════════════════════════════════════════════════════════════════════

def slide_extended_section_break(prs):
    sl = blank_slide(prs); fill_bg(sl)
    add_rect(sl, Inches(0), Inches(2.4), SLIDE_W, Inches(2.4), fill=RGBColor(0x1A, 0x23, 0x3D))
    add_text(sl, "§4.9  Extended Analyses",
             Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.1),
             size=36, bold=True, colour=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(sl,
             "Semantic similarity  ·  Bootstrap CI  ·  SVD rank  ·  "
             "Taxonomy grouping  ·  Quality regression  ·  Neutral check  ·  Network graph",
             Inches(0.5), Inches(3.78), Inches(12.3), Inches(0.5),
             size=14, colour=C_GREY, align=PP_ALIGN.CENTER)
    add_text(sl, "All analyses use jury-confirmed rows only  (N = 273)",
             Inches(1.5), Inches(6.35), Inches(10.3), Inches(0.4),
             size=13, colour=C_YELLOW, align=PP_ALIGN.CENTER)

    analyses = [
        (C_ACCENT,  "Latent Priming",  "Spearman ρ(s[β,γ], M[β,γ])"),
        (C_GREEN,   "Bootstrap CI",    "95% cell stability intervals"),
        (C_YELLOW,  "SVD Rank",        "Intrinsic dimensionality of M"),
        (C_RED,     "Taxonomy Test",   "Within- vs. across-group Δ"),
        (C_GREY,    "Quality Reg.",    "Does stronger signal → higher Δ?"),
        (C_ACCENT,  "Neutral Check",   "Adversarial baseline validation"),
        (C_GREEN,   "Network Graph",   "Directed induction topology"),
    ]
    for i, (col, title, sub) in enumerate(analyses):
        x = Inches(0.38 + i * 1.88)
        add_rect(sl, x, Inches(4.55), Inches(1.72), Inches(1.42),
                 fill=RGBColor(0x1E, 0x29, 0x4A))
        add_rect(sl, x, Inches(4.55), Inches(1.72), Inches(0.3), fill=col)
        add_text(sl, title, x + Inches(0.06), Inches(4.57), Inches(1.6), Inches(0.26),
                 size=9.5, bold=True, colour=C_BG, align=PP_ALIGN.CENTER)
        add_text(sl, sub, x + Inches(0.06), Inches(4.9), Inches(1.6), Inches(0.9),
                 size=9, colour=C_WHITE, align=PP_ALIGN.CENTER)


def slide_latent_priming_deep(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.9a  Latent Priming — Deep Evidence",
              "Semantic similarity correlation  (Panel C)  &  SVD rank decomposition of M")
    section_tag(sl, "§4.9a")

    sem_img = _try_img(_EXT_PLOT_DIR, "semantic_similarity.png")
    svd_img = _try_img(_EXT_PLOT_DIR, "svd_rank.png")

    if sem_img:
        add_image_file(sl, sem_img, Inches(0.25), Inches(1.2), Inches(6.6), Inches(3.85))
    else:
        _placeholder_box(sl, Inches(0.25), Inches(1.2), Inches(6.6), Inches(3.85),
                         "semantic_similarity.png\n(run: python extended_analysis.py)")

    if svd_img:
        add_image_file(sl, svd_img, Inches(6.95), Inches(1.2), Inches(6.1), Inches(3.85))
    else:
        _placeholder_box(sl, Inches(6.95), Inches(1.2), Inches(6.1), Inches(3.85),
                         "svd_rank.png\n(run: python extended_analysis.py)")

    add_rect(sl, Inches(0.25), Inches(5.2), Inches(12.85), Inches(2.05),
             fill=RGBColor(0x1E, 0x29, 0x4A))
    add_text(sl, "Semantic Similarity Test  (Latent Priming — Panel C):",
             Inches(0.4), Inches(5.28), Inches(6.1), Inches(0.35),
             size=11, bold=True, colour=C_ACCENT)
    add_text(sl,
             "Spearman ρ between cosine similarity s[β,γ] and coupling M[β,γ], off-diagonal only\n"
             "Mantel permutation test (2 000 perms) guards against spatial autocorrelation\n"
             "→ ρ > 0 & p < 0.05: semantically closer biases induce each other more — latent priming confirmed",
             Inches(0.4), Inches(5.65), Inches(6.1), Inches(1.5),
             size=10, colour=C_WHITE)
    add_text(sl, "SVD Rank Decomposition  M = UΣVᵀ:",
             Inches(6.7), Inches(5.28), Inches(6.0), Inches(0.35),
             size=11, bold=True, colour=C_ACCENT)
    add_text(sl,
             "Low effective rank → few dominant induction patterns (structured mechanism)\n"
             "High rank → noisy scenario-specific effects\n"
             "U₁ / V₁ loading bars identify which biases drive the primary induction direction\n"
             "Cumulative variance plot: how many components explain ≥ 90% of M variance?",
             Inches(6.7), Inches(5.65), Inches(6.0), Inches(1.5),
             size=10, colour=C_WHITE)


def slide_bootstrap_taxonomy(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.9b  Matrix Stability & Taxonomy Structure",
              "Bootstrap 95% CI width  &  Within- vs. across-group induction test")
    section_tag(sl, "§4.9b")

    boot_img = _try_img(_EXT_PLOT_DIR, "bootstrap_ci.png")
    tax_img  = _try_img(_EXT_PLOT_DIR, "taxonomy_test.png")

    if boot_img:
        add_image_file(sl, boot_img, Inches(0.25), Inches(1.2), Inches(6.6), Inches(3.85))
    else:
        _placeholder_box(sl, Inches(0.25), Inches(1.2), Inches(6.6), Inches(3.85),
                         "bootstrap_ci.png\n(run: python extended_analysis.py)")

    if tax_img:
        add_image_file(sl, tax_img, Inches(6.95), Inches(1.2), Inches(6.1), Inches(3.85))
    else:
        _placeholder_box(sl, Inches(6.95), Inches(1.2), Inches(6.1), Inches(3.85),
                         "taxonomy_test.png\n(run: python extended_analysis.py)")

    add_rect(sl, Inches(0.25), Inches(5.2), Inches(12.85), Inches(2.05),
             fill=RGBColor(0x1E, 0x29, 0x4A))
    add_text(sl, "Bootstrap 95% CI  (2 000 resamples per cell):",
             Inches(0.4), Inches(5.28), Inches(6.1), Inches(0.35),
             size=11, bold=True, colour=C_ACCENT)
    add_text(sl,
             "Left heatmap: point estimates (identical to M[β,γ])\n"
             "Right heatmap: CI width — wider = less stable, more data needed\n"
             "→ Flags cells where n is too small to trust the estimate\n"
             "Guides where next pilot scenarios should be allocated",
             Inches(0.4), Inches(5.65), Inches(6.1), Inches(1.5),
             size=10, colour=C_WHITE)
    add_text(sl, "Taxonomy Grouping Test  (permutation, 5 000 perms):",
             Inches(6.7), Inches(5.28), Inches(6.0), Inches(0.35),
             size=11, bold=True, colour=C_ACCENT)
    add_text(sl,
             "Groups: Probability/Estimation · Social/Conformity · Temporal/Planning · Value/Risk\n"
             "H₀: Within-group mean Δ = Across-group mean Δ\n"
             "Violin: within / across / same-bias (diagonal) distributions\n"
             "→ If within > across: biases preferentially induce within their cognitive family",
             Inches(6.7), Inches(5.65), Inches(6.0), Inches(1.5),
             size=10, colour=C_WHITE)


def slide_regression_network(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.9c  Quality Regression, Neutral Check & Induction Network",
              "Does message strength predict Δ?  ·  Baseline validation  ·  Induction topology")
    section_tag(sl, "§4.9c")

    reg_img = _try_img(_EXT_PLOT_DIR, "message_quality_regression.png")
    neu_img = _try_img(_EXT_PLOT_DIR, "neutral_check.png")
    net_img = _try_img(_EXT_PLOT_DIR, "network_graph.png")

    if reg_img:
        add_image_file(sl, reg_img, Inches(0.25), Inches(1.2), Inches(6.6), Inches(2.75))
    else:
        _placeholder_box(sl, Inches(0.25), Inches(1.2), Inches(6.6), Inches(2.75),
                         "message_quality_regression.png")

    if neu_img:
        add_image_file(sl, neu_img, Inches(6.95), Inches(1.2), Inches(6.1), Inches(2.75))
    else:
        _placeholder_box(sl, Inches(6.95), Inches(1.2), Inches(6.1), Inches(2.75),
                         "neutral_check.png")

    if net_img:
        add_image_file(sl, net_img, Inches(3.8), Inches(4.1), Inches(5.75), Inches(3.15))
    else:
        _placeholder_box(sl, Inches(0.25), Inches(4.1), Inches(12.85), Inches(3.15),
                         "network_graph.png\n(run: python extended_analysis.py)")

    add_text(sl,
             "OLS: Δ ~ jury_target_rating + purity + naturalness + sim_assertiveness + "
             "sim_emotional_intensity + sim_sentiment + sim_word_count + bias FEs  (HC3 robust SE)",
             Inches(0.25), Inches(4.0), Inches(6.5), Inches(0.42),
             size=8.5, colour=C_GREY)
    add_text(sl,
             "Neutral baseline t-test vs. zero  (contamination check — p < 0.05 = problem)",
             Inches(6.95), Inches(4.0), Inches(5.9), Inches(0.42),
             size=8.5, colour=C_GREY)
    add_text(sl,
             "Network: directed edge from human bias → target bias when |mean Δ| ≥ 0.08  "
             "·  green = amplification  ·  red = suppression  ·  node colour = taxonomy group",
             Inches(0.25), Inches(7.2), Inches(12.85), Inches(0.38),
             size=9, colour=C_GREY, align=PP_ALIGN.CENTER)


def slide_multi_llm_plan(prs):
    sl = blank_slide(prs); fill_bg(sl)
    title_bar(sl, "§4.10  Multi-LLM Replication Study",
              "Same protocol across 5 models — tests generalisability of the coupling structure")
    section_tag(sl, "§4.10")

    models = [
        ("GPT-4o",            "OpenAI",    "88.0", "Baseline  (current results)"),
        ("GPT-4o-Mini",       "OpenAI",    "82.0", "Smaller / cheaper — same family"),
        ("Claude-3.5-Sonnet", "Anthropic", "88.3", "RLHF-heavy, different alignment"),
        ("Llama-3.1-70B",     "Meta",      "83.6", "Open weights, SFT only"),
        ("Gemini-1.5-Flash",  "Google",    "78.9", "Multimodal, different pretraining"),
    ]
    headers      = ["Model", "Provider", "MMLU", "Role in study"]
    col_w        = [2.5, 1.4, 0.85, 4.3]
    col_x        = [0.35]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    add_rect(sl, Inches(0.35), Inches(1.25), Inches(9.1), Inches(0.42), fill=C_ACCENT)
    for h, w, x in zip(headers, col_w, col_x):
        add_text(sl, h, Inches(x + 0.08), Inches(1.27), Inches(w - 0.1), Inches(0.36),
                 size=11, bold=True, colour=C_BG)

    for ri, (model, provider, mmlu, role) in enumerate(models):
        bg = RGBColor(0x1A, 0x23, 0x3D) if ri % 2 == 0 else RGBColor(0x15, 0x1E, 0x35)
        add_rect(sl, Inches(0.35), Inches(1.72 + ri * 0.52), Inches(9.1), Inches(0.50), fill=bg)
        for val, w, x, col in zip([model, provider, mmlu, role], col_w, col_x,
                                   [C_ACCENT, C_WHITE, C_WHITE, C_WHITE]):
            add_text(sl, val, Inches(x + 0.08), Inches(1.75 + ri * 0.52),
                     Inches(w - 0.1), Inches(0.42), size=10.5, colour=col)

    add_rect(sl, Inches(9.65), Inches(1.25), Inches(3.45), Inches(5.85),
             fill=RGBColor(0x1A, 0x23, 0x3D))
    add_rect(sl, Inches(9.65), Inches(1.25), Inches(3.45), Inches(0.42), fill=C_GREEN)
    add_text(sl, "Planned analyses", Inches(9.72), Inches(1.27), Inches(3.3), Inches(0.36),
             size=11, bold=True, colour=C_BG)
    plan_items = [
        "Coupling matrix per model",
        "Replication map  (fraction consistent)",
        "Susceptibility radar chart",
        "Capability vs. susceptibility scatter",
        "Cross-model Spearman ρ of M",
        "Model×bias LMM  (fixed effects)",
    ]
    for i, item in enumerate(plan_items):
        add_text(sl, f"▸  {item}",
                 Inches(9.72), Inches(1.75 + i * 0.57), Inches(3.3), Inches(0.52),
                 size=10.5, colour=C_WHITE)

    add_rect(sl, Inches(0.35), Inches(6.58), Inches(12.75), Inches(0.72),
             fill=RGBColor(0x1E, 0x29, 0x4A))
    add_text(sl,
             "H-A: In-Group suppression replicates across all 5 models  (most robust signal)\n"
             "H-B: Coupling matrix ρ(GPT-4o, GPT-4o-Mini) > ρ(GPT-4o, Llama)  "
             "— model family matters\n"
             "H-C: Higher MMLU ≠ lower bias susceptibility  "
             "(capability and bias are orthogonal dimensions)",
             Inches(0.5), Inches(6.62), Inches(12.5), Inches(0.65),
             size=9.5, colour=C_YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    print("Loading data …")
    jp, pivot = load()

    print("Building presentation …")
    prs = new_prs()

    slide_title(prs)
    slide_agenda(prs)
    slide_hypothesis(prs)
    slide_three_mechanisms(prs)
    slide_protocol(prs)
    slide_jury(prs)
    slide_metric(prs)
    slide_matrix(prs, jp, pivot)
    slide_mechanism_tests(prs, jp, pivot)
    slide_row_col(prs, jp)
    slide_significance(prs, jp)
    slide_causal(prs, jp)
    slide_summary(prs)
    # Extended analyses  §4.9 — §4.10
    slide_extended_section_break(prs)
    slide_latent_priming_deep(prs)
    slide_bootstrap_taxonomy(prs)
    slide_regression_network(prs)
    slide_multi_llm_plan(prs)

    os.makedirs(os.path.dirname(_OUT_PPTX), exist_ok=True)
    prs.save(_OUT_PPTX)
    print(f"\nSaved → {_OUT_PPTX}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
